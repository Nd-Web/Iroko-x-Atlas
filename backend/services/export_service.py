"""
Export Service — one-click PDF and PPTX generation from AI answers.

Design language: MTN Nigeria brand identity
  Yellow  #FFCB00  — primary accent, headers, highlights
  Black   #1A1A1A  — headings and dark backgrounds
  Dark    #2D2D2D  — body text
  Gray    #6B6B6B  — captions, metadata
  White   #FFFFFF  — slide / page backgrounds

Fonts:
  PDF  — Open Sans (Liberation Sans fallback → Helvetica)
  PPTX — Calibri (Office default; metrically modern and clean)
"""
import io
import json
import logging
import os
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# ── MTN Brand palette ─────────────────────────────────────────────────────────
MTN_YELLOW = "#FFCB00"
MTN_BLACK  = "#1A1A1A"
MTN_DARK   = "#2D2D2D"
MTN_GRAY   = "#6B6B6B"
MTN_LGRAY  = "#F5F5F5"

_Y  = (0xFF, 0xCB, 0x00)
_BK = (0x1A, 0x1A, 0x1A)
_DK = (0x2D, 0x2D, 0x2D)
_GR = (0x6B, 0x6B, 0x6B)
_WH = (0xFF, 0xFF, 0xFF)
_LG = (0xF5, 0xF5, 0xF5)

FONT_BODY = "Calibri"
FONT_BOLD = "Calibri"

# ── Asset helpers ─────────────────────────────────────────────────────────────

def _logo_path() -> str | None:
    """Return the absolute path to the MTN logo, or None if not found."""
    candidates = [
        Path(__file__).parent.parent / "assets" / "mtn_logo.jpg",
        Path(__file__).parent.parent / "assets" / "mtn_logo.png",
        Path("/app/assets/mtn_logo.jpg"),
        Path("/app/assets/mtn_logo.png"),
    ]
    for p in candidates:
        if p.exists():
            return str(p)
    return None


def _pdf_font_name(bold: bool = False) -> str:
    """Best available font name for reportlab — Open Sans > Liberation > Helvetica."""
    _OPEN_SANS_BOLD   = "/usr/share/fonts/truetype/open-sans/OpenSans-Bold.ttf"
    _OPEN_SANS_REG    = "/usr/share/fonts/truetype/open-sans/OpenSans-Regular.ttf"
    _LIB_BOLD         = "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf"
    _LIB_REG          = "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf"

    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        if bold:
            for path, name in [(_OPEN_SANS_BOLD, "OpenSans-Bold"), (_LIB_BOLD, "LiberationSans-Bold")]:
                if os.path.exists(path):
                    try:
                        pdfmetrics.registerFont(TTFont(name, path))
                        return name
                    except Exception:
                        pass
            return "Helvetica-Bold"
        else:
            for path, name in [(_OPEN_SANS_REG, "OpenSans"), (_LIB_REG, "LiberationSans")]:
                if os.path.exists(path):
                    try:
                        pdfmetrics.registerFont(TTFont(name, path))
                        return name
                    except Exception:
                        pass
            return "Helvetica"
    except Exception:
        return "Helvetica-Bold" if bold else "Helvetica"


# ─────────────────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────────────────

def generate_pdf(title: str, answer: str, citations: list, map_data: list = [], fraud_data: dict | None = None) -> bytes:
    """Render a branded MTN Nigeria A4 PDF from a markdown answer."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Image as RLImage

    YELLOW = HexColor(MTN_YELLOW)
    BLACK  = HexColor(MTN_BLACK)
    DARK   = HexColor(MTN_DARK)
    GRAY   = HexColor(MTN_GRAY)
    LGRAY  = HexColor(MTN_LGRAY)

    font_reg  = _pdf_font_name(bold=False)
    font_bold = _pdf_font_name(bold=True)

    buf      = io.BytesIO()
    logo     = _logo_path()

    # ── Page callback — draws header + footer on every page ──────────────────
    def _draw_page(canv, doc):
        W, H = A4
        canv.saveState()

        # Yellow top bar
        canv.setFillColor(YELLOW)
        canv.rect(0, H - 20*mm, W, 20*mm, fill=1, stroke=0)

        # Logo in top-left of header
        if logo:
            try:
                canv.drawImage(logo, 8*mm, H - 17*mm, width=22*mm, height=14*mm,
                               preserveAspectRatio=True, anchor="sw", mask="auto")
            except Exception:
                canv.setFillColor(BLACK)
                canv.setFont(font_bold, 11)
                canv.drawString(22*mm, H - 11*mm, "MTN Nigeria")
        else:
            canv.setFillColor(BLACK)
            canv.setFont(font_bold, 11)
            canv.drawString(22*mm, H - 11*mm, "MTN Nigeria")

        canv.setFillColor(BLACK)
        canv.setFont(font_reg, 8)
        canv.drawRightString(W - 12*mm, H - 10*mm, "Powered by Iroko AI")

        # Thin black separator below yellow bar
        canv.setFillColor(BLACK)
        canv.rect(0, H - 21*mm, W, 1*mm, fill=1, stroke=0)

        # Black bottom bar
        canv.rect(0, 0, W, 11*mm, fill=1, stroke=0)
        canv.setFillColor(YELLOW)
        canv.setFont(font_reg, 7)
        canv.drawString(12*mm, 3.8*mm, "Confidential  ·  Internal Use Only  ·  MTN Nigeria Communications Plc")
        canv.setFillColor(HexColor("#AAAAAA"))
        canv.drawRightString(W - 12*mm, 3.8*mm, f"Page {doc.page}")
        canv.restoreState()

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=22*mm, rightMargin=22*mm,
        topMargin=26*mm, bottomMargin=18*mm,
    )

    base = getSampleStyleSheet()
    def S(name, **kw):
        return ParagraphStyle(name, parent=base["Normal"], **kw)

    title_st = S("pt",  fontSize=22, textColor=BLACK, fontName=font_bold,
                 spaceAfter=4, leading=28)
    h2       = S("ph2", fontSize=13, textColor=BLACK, fontName=font_bold,
                 spaceAfter=4, spaceBefore=16, leading=17)
    h3       = S("ph3", fontSize=11, textColor=DARK,  fontName=font_bold,
                 spaceAfter=3, spaceBefore=10)
    body     = S("pb",  fontSize=10, textColor=DARK,  fontName=font_reg,
                 leading=17, spaceAfter=4)
    bul      = S("pbu", fontSize=10, textColor=DARK,  fontName=font_reg,
                 leading=17, spaceAfter=3, leftIndent=14)
    cap      = S("pc",  fontSize=8,  textColor=GRAY,  fontName=font_reg,
                 leading=12)

    badge_st = S("pbadge", fontSize=8, textColor=YELLOW, fontName=font_bold,
                 leading=10, spaceAfter=3, backColor=BLACK,
                 leftIndent=0, borderPadding=(3, 6, 3, 6))
    meta_st  = S("pmeta", fontSize=8, textColor=GRAY, fontName=font_reg,
                 leading=12, spaceAfter=6)
    title_st = S("pt", fontSize=26, textColor=BLACK, fontName=font_bold,
                 spaceAfter=6, leading=32)

    from datetime import date as _date
    _today = _date.today().strftime("%d %B %Y")

    story = []
    story.append(Spacer(1, 4*mm))
    story.append(Paragraph("INTELLIGENCE REPORT", badge_st))
    story.append(Spacer(1, 2*mm))
    story.append(Paragraph(_esc(title), title_st))
    story.append(HRFlowable(width="100%", thickness=3, color=YELLOW, spaceAfter=4))
    story.append(Paragraph(
        f"MTN Nigeria Communications Plc  ·  Confidential — Internal Use Only  ·  {_today}",
        meta_st))
    story.append(Spacer(1, 4*mm))

    _md_to_story(answer, story, h2, h3, body, bul, YELLOW, LGRAY)

    if map_data:
        try:
            img_bytes = _generate_map_image(map_data)
            story.append(Spacer(1, 6*mm))
            story.append(Paragraph("Network Status Map — Nigeria", h2))
            story.append(HRFlowable(width="100%", thickness=1, color=YELLOW, spaceAfter=4))
            story.append(Spacer(1, 2*mm))
            story.append(RLImage(io.BytesIO(img_bytes), width=170*mm, height=110*mm))
            story.append(Spacer(1, 4*mm))
        except Exception as _e:
            logger.warning(f"Map image in PDF skipped: {_e}")

    if fraud_data:
        try:
            story.append(Spacer(1, 6*mm))
            story.append(Paragraph("Fraud Intelligence Summary", h2))
            story.append(HRFlowable(width="100%", thickness=1, color=YELLOW, spaceAfter=4))
            story.append(Spacer(1, 2*mm))
            hi  = fraud_data.get("high_risk", 0)
            med = fraud_data.get("medium_risk", 0)
            exp = fraud_data.get("total_exposure_ngn", 0)
            story.append(Paragraph(
                f"<b>{hi} HIGH</b>  ·  {med} MEDIUM risk signals active  ·  "
                f"Financial exposure: <b>₦{exp:,.0f}</b>", body))
            story.append(Spacer(1, 3*mm))
            for s in (fraud_data.get("signals") or [])[:6]:
                risk_label = s.get("risk", "MEDIUM")
                story.append(Paragraph(
                    f"<b>[{risk_label}]  {_esc(s.get('title', ''))} — {_esc(s.get('id', ''))}</b>", body))
                story.append(Paragraph(_esc(s.get("detail", "")[:240]), cap))
                if s.get("amount_ngn"):
                    story.append(Paragraph(
                        f"Exposure: ₦{s['amount_ngn']:,.0f}  ·  Status: {s.get('status', '—')}", cap))
                story.append(Spacer(1, 2*mm))
        except Exception as _e:
            logger.warning(f"Fraud section in PDF skipped: {_e}")

    if citations:
        story.append(Spacer(1, 6*mm))
        story.append(HRFlowable(width="100%", thickness=1, color=YELLOW, spaceAfter=4))
        story.append(Paragraph("Sources &amp; References", h2))
        for i, c in enumerate(citations, 1):
            if isinstance(c, str):
                story.append(Paragraph(f"[{i}]  <b>{_esc(c)}</b>", body))
                story.append(Spacer(1, 2*mm))
                continue
            doc_title = c.get("document_title") or c.get("source") or c.get("document_id", "Unknown")
            excerpt   = c.get("excerpt", "")
            story.append(Paragraph(f"[{i}]  <b>{_esc(doc_title)}</b>", body))
            if excerpt:
                snip = excerpt[:180] + ("…" if len(excerpt) > 180 else "")
                story.append(Paragraph(_esc(snip), cap))
            story.append(Spacer(1, 2*mm))

    doc.build(story, onFirstPage=_draw_page, onLaterPages=_draw_page)
    return buf.getvalue()


# ── Markdown → reportlab ──────────────────────────────────────────────────────

def _md_to_story(text, story, h2, h3, body, bul, YELLOW, LGRAY):
    from reportlab.platypus import Spacer, Paragraph
    from reportlab.lib.units import mm

    lines = text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()

        if line.startswith("### "):
            story.append(Paragraph(_inline(_esc(line[4:])), h3))
        elif line.startswith("## "):
            story.append(Paragraph(_inline(_esc(line[3:])), h2))
        elif line.startswith("# "):
            story.append(Paragraph(_inline(_esc(line[2:])), h2))
        elif line.startswith("- ") or line.startswith("* "):
            story.append(Paragraph("▸  " + _inline(_esc(line[2:])), bul))
        elif re.match(r"^\d+\.\s", line):
            content = re.sub(r"^\d+\.\s", "", line)
            story.append(Paragraph("▸  " + _inline(_esc(content)), bul))
        elif line.startswith("| "):
            rows_raw = [line]
            j = i + 1
            while j < len(lines) and lines[j].lstrip().startswith("|"):
                rows_raw.append(lines[j].rstrip())
                j += 1
            i = j - 1
            _table_to_story(rows_raw, story, YELLOW, LGRAY)
        elif line in ("", "---", "***"):
            story.append(Spacer(1, 3*mm))
        else:
            if line.strip():
                story.append(Paragraph(_inline(_esc(line)), body))
        i += 1


def _inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
    text = re.sub(r"\*(.+?)\*",     r"<i>\1</i>", text)
    text = re.sub(r"`(.+?)`",       r"<font name='Courier'>\1</font>", text)
    return text


def _esc(text: str) -> str:
    return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _table_to_story(rows_raw, story, YELLOW, LGRAY):
    from reportlab.platypus import Table, TableStyle, Spacer, Paragraph
    from reportlab.lib.colors import HexColor, white
    from reportlab.lib.units import cm, mm
    from reportlab.lib.styles import getSampleStyleSheet

    st = getSampleStyleSheet()["Normal"]
    st.fontSize = 9

    rows = []
    for raw in rows_raw:
        if re.match(r"^\|[-:| ]+\|$", raw.strip()):
            continue
        cells = [c.strip() for c in raw.strip("|").split("|")]
        rows.append(cells)

    if not rows:
        return

    data    = [[Paragraph(_inline(_esc(c)), st) for c in row] for row in rows]
    col_cnt = max(len(r) for r in data)
    col_w   = (17 * cm) / col_cnt

    t = Table(data, colWidths=[col_w] * col_cnt, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND",     (0, 0), (-1,  0), YELLOW),
        ("TEXTCOLOR",      (0, 0), (-1,  0), HexColor(MTN_BLACK)),
        ("FONTNAME",       (0, 0), (-1,  0), "Helvetica-Bold"),
        ("FONTSIZE",       (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, LGRAY]),
        ("GRID",           (0, 0), (-1, -1), 0.3, HexColor("#DDDDDD")),
        ("ALIGN",          (0, 0), (-1, -1), "LEFT"),
        ("VALIGN",         (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING",     (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING",  (0, 0), (-1, -1), 5),
        ("LEFTPADDING",    (0, 0), (-1, -1), 6),
    ]))
    story.append(t)
    story.append(Spacer(1, 4*mm))


# ─────────────────────────────────────────────────────────────────────────────
# Network heatmap image — OSM tiles via staticmap (same source as Leaflet)
# ─────────────────────────────────────────────────────────────────────────────

def _generate_map_image(map_data: list) -> bytes:
    """Render Nigeria network status map using real OSM tiles (staticmap).

    Uses the same tile source as the Leaflet frontend, so the map boundary,
    roads, and city labels are all accurate and complete.  Circle markers are
    drawn on top via Pillow; pixel positions are derived from the Web Mercator
    projection at the fixed zoom level.
    """
    import math
    from staticmap import StaticMap, CircleMarker as _SMCircle
    from PIL import ImageDraw, ImageFont, Image as _PILImage  # staticmap depends on Pillow

    # Nigeria viewport — fixed centre and zoom so projections are deterministic
    _ZOOM       = 6
    _CENTER_LON = 8.3   # slightly east of geographic centre to balance the shape
    _CENTER_LAT = 9.1
    _IMG_W      = 1600
    _IMG_H      = 1100

    _STATUS_RGB: dict[str, tuple[int, int, int]] = {
        "operational": (22,  163, 74),
        "degraded":    (217, 119,  6),
        "down":        (220,  38, 38),
    }

    # ── Fetch OSM tiles ───────────────────────────────────────────────────────
    m = StaticMap(
        _IMG_W, _IMG_H,
        url_template="https://tile.openstreetmap.org/{z}/{x}/{y}.png",
        padding_x=0, padding_y=0,
        headers={"User-Agent": "IrokoAI-ExportService/1.0 (+https://iroko.ai)"},
    )
    img = m.render(zoom=_ZOOM, center=(_CENTER_LON, _CENTER_LAT))

    # ── Pixel-coordinate helper (Web Mercator → image pixels) ─────────────────
    _tile_px = 256 * (2 ** _ZOOM)  # total world width/height in pixels at this zoom

    def _to_px(lat: float, lon: float) -> tuple[int, int]:
        cx = (_CENTER_LON + 180) / 360 * _tile_px
        cr = math.radians(_CENTER_LAT)
        cy = (1 - math.log(math.tan(cr) + 1 / math.cos(cr)) / math.pi) / 2 * _tile_px
        x  = (lon + 180) / 360 * _tile_px
        r  = math.radians(lat)
        y  = (1 - math.log(math.tan(r)  + 1 / math.cos(r))  / math.pi) / 2 * _tile_px
        return (
            int(_IMG_W / 2 + (x - cx)),
            int(_IMG_H / 2 + (y - cy)),
        )

    # ── Load fonts (best-effort) ───────────────────────────────────────────────
    _FONT_PATHS_BOLD = [
        "/usr/share/fonts/truetype/open-sans/OpenSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    ]
    _FONT_PATHS_REG = [
        "/usr/share/fonts/truetype/open-sans/OpenSans-Regular.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]

    def _load_font(paths: list[str], size: int):
        for p in paths:
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                pass
        return ImageFont.load_default()

    font_label = _load_font(_FONT_PATHS_BOLD, 17)
    font_avail = _load_font(_FONT_PATHS_REG,  14)
    font_leg   = _load_font(_FONT_PATHS_BOLD, 15)

    draw = ImageDraw.Draw(img, "RGBA")

    # ── Draw markers ──────────────────────────────────────────────────────────
    for r in map_data:
        px, py   = _to_px(r["latitude"], r["longitude"])
        rgb      = _STATUS_RGB.get(r.get("status", "operational"), (22, 163, 74))
        radius   = max(14, min(38, int(r.get("site_count", 1) * 3.5)))

        # White halo → coloured circle
        draw.ellipse(
            [px - radius - 3, py - radius - 3, px + radius + 3, py + radius + 3],
            fill=(255, 255, 255, 230),
        )
        draw.ellipse(
            [px - radius, py - radius, px + radius, py + radius],
            fill=(*rgb, 215),
        )

        # Region name — white pill above the circle
        label = r.get("region", "")
        bb    = draw.textbbox((0, 0), label, font=font_label)
        tw, th = bb[2] - bb[0], bb[3] - bb[1]
        lx = px - tw // 2
        ly = py - radius - th - 14
        pad = 4
        draw.rounded_rectangle(
            [lx - pad, ly - pad, lx + tw + pad, ly + th + pad],
            radius=5, fill=(255, 255, 255, 225),
        )
        draw.text((lx, ly), label, font=font_label, fill=(26, 26, 26))

        # Availability % — below the circle
        avail = r.get("availability_pct")
        if avail is not None:
            txt  = f"{avail}%"
            bb2  = draw.textbbox((0, 0), txt, font=font_avail)
            tw2  = bb2[2] - bb2[0]
            a_color = (220, 38, 38) if avail < 95 else (217, 119, 6) if avail < 99 else (22, 163, 74)
            draw.text((px - tw2 // 2, py + radius + 6), txt, font=font_avail, fill=a_color)

    # ── Legend (bottom-right) ─────────────────────────────────────────────────
    leg_items = [
        ("Operational", (22, 163, 74)),
        ("Degraded",    (217, 119, 6)),
        ("Down",        (220, 38, 38)),
    ]
    leg_row_h = 34
    leg_w, leg_h = 195, len(leg_items) * leg_row_h + 20
    lx0 = _IMG_W - leg_w - 18
    ly0 = _IMG_H - leg_h - 18
    draw.rounded_rectangle(
        [lx0 - 10, ly0 - 8, lx0 + leg_w, ly0 + leg_h],
        radius=8, fill=(255, 255, 255, 230),
    )
    for i, (lbl, rgb) in enumerate(leg_items):
        cy_leg = ly0 + 10 + i * leg_row_h
        draw.ellipse([lx0, cy_leg, lx0 + 18, cy_leg + 18], fill=rgb)
        draw.text((lx0 + 26, cy_leg + 1), lbl, font=font_leg, fill=(45, 45, 45))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# DALL-E / GPT-Image hero generation (optional — requires image deployment)
# ─────────────────────────────────────────────────────────────────────────────

async def _generate_dalle_image(prompt: str, size: str = "1792x1024") -> bytes | None:
    """
    Call Azure OpenAI image generation (DALL-E 3 / gpt-image-1).
    Returns PNG bytes, or None if no deployment is configured or the call fails.
    Set AZURE_OPENAI_IMAGE_DEPLOYMENT env var to enable (e.g. 'dall-e-3').
    """
    deployment = os.getenv("AZURE_OPENAI_IMAGE_DEPLOYMENT")
    if not deployment:
        return None
    try:
        import httpx
        endpoint = os.getenv("AZURE_OPENAI_ENDPOINT", "").rstrip("/")
        api_key  = os.getenv("AZURE_OPENAI_API_KEY", "")
        api_ver  = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-01")
        url = f"{endpoint}/openai/deployments/{deployment}/images/generations?api-version={api_ver}"
        payload = {"prompt": prompt, "n": 1, "size": size, "quality": "standard"}
        async with httpx.AsyncClient(timeout=45) as client:
            r = await client.post(url, json=payload, headers={"api-key": api_key})
            r.raise_for_status()
            img_url = r.json()["data"][0]["url"]
            img_resp = await client.get(img_url, timeout=30)
            img_resp.raise_for_status()
            return img_resp.content
    except Exception as e:
        logger.warning(f"DALL-E image generation skipped: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────

async def generate_pptx(title: str, answer: str, citations: list, map_data: list = [], fraud_data: dict | None = None) -> bytes:
    """
    Generate a 16:9 PowerPoint deck with full MTN Nigeria branding.

    Slides:
      1  Title          — black bg, MTN logo, DALL-E hero image (if configured)
      2  Agenda         — clean white, numbered section index
      3+ Content        — white bg, yellow header strip, Calibri bullets
      N  Key Takeaways  — yellow bg, high-impact closer
      N+1 Fraud         — dark bg (if fraud_data)
      N+2 Network Map   — dark bg (if map_data)
      N+3 Sources       — light gray bg
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    slide_data = await _structure_slides(title, answer)
    slides     = slide_data.get("slides", [])

    Y  = RGBColor(*_Y)
    BK = RGBColor(*_BK)
    DK = RGBColor(*_DK)
    GR = RGBColor(*_GR)
    WH = RGBColor(*_WH)
    LG = RGBColor(*_LG)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    logo = _logo_path()

    # Generate hero image for title slide (optional)
    hero_prompt = (
        f"Professional, photorealistic wide-angle illustration for a corporate "
        f"intelligence report titled '{title[:80]}'. "
        f"Modern Nigerian telecom infrastructure, aerial city view, clean and minimal. "
        f"No text. Cinematic lighting."
    )
    hero_bytes = await _generate_dalle_image(hero_prompt)

    # ── 1. Title slide ────────────────────────────────────────────────────────
    ts = prs.slides.add_slide(blank)
    _bg(ts, BK)

    # Hero background image (left 2/3 of slide)
    if hero_bytes:
        try:
            ts.shapes.add_picture(io.BytesIO(hero_bytes),
                                  Inches(0), Inches(0), Inches(8.5), Inches(7.5))
            # dark overlay to keep text readable
            _rect(ts, Inches(0), Inches(0), Inches(8.5), Inches(7.5),
                  RGBColor(0x0A, 0x0A, 0x0A))
            # reduce overlay opacity via XML directly
            ov = ts.shapes[-1]
            sp_el = ov._element
            sp_pr = sp_el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
            if sp_pr is None:
                pass  # fallback — the rect is already dark
        except Exception:
            pass

    # Right panel — dark yellow strip
    _rect(ts, Inches(9.0), Inches(0), Inches(4.33), Inches(7.5), RGBColor(0xFF, 0xCB, 0x00))
    _rect(ts, Inches(8.85), Inches(0), Inches(0.15), Inches(7.5), RGBColor(0x11, 0x11, 0x11))

    # Logo top-right
    if logo:
        try:
            ts.shapes.add_picture(logo, Inches(9.4), Inches(0.3), Inches(2.9), Inches(1.5))
        except Exception:
            _tb(ts, Inches(9.4), Inches(0.4), Inches(3.5), Inches(0.8),
                "MTN", Pt(32), BK, bold=True, align=PP_ALIGN.CENTER)
    else:
        _tb(ts, Inches(9.4), Inches(0.4), Inches(3.5), Inches(0.8),
            "MTN", Pt(32), BK, bold=True, align=PP_ALIGN.CENTER)
        _tb(ts, Inches(9.4), Inches(1.2), Inches(3.5), Inches(0.4),
            "NIGERIA", Pt(11), BK, align=PP_ALIGN.CENTER)

    # Category badge + title text (over hero or dark bg)
    _rect(ts, Inches(0.55), Inches(1.85), Inches(2.4), Inches(0.35), Y)
    _tb(ts, Inches(0.58), Inches(1.87), Inches(2.35), Inches(0.32),
        "INTELLIGENCE REPORT", Pt(9), BK, bold=True)
    _tb(ts, Inches(0.55), Inches(2.35), Inches(8.0), Inches(2.5),
        title, Pt(34), WH, bold=True)
    _rect(ts, Inches(0.55), Inches(5.0), Inches(4.5), Inches(0.05), Y)
    from datetime import date as _dt
    _slide_date = _dt.today().strftime("%d %B %Y")
    _tb(ts, Inches(0.55), Inches(5.15), Inches(8.0), Inches(0.45),
        f"MTN Nigeria  ·  Confidential  ·  {_slide_date}  ·  Powered by Iroko AI", Pt(11), GR)

    # Right-panel tagline
    _tb(ts, Inches(9.1), Inches(5.5), Inches(3.9), Inches(0.5),
        "Enterprise Intelligence", Pt(11), BK, bold=True, align=PP_ALIGN.CENTER)
    _tb(ts, Inches(9.1), Inches(5.95), Inches(3.9), Inches(0.4),
        "Confidential · Internal Only", Pt(9), RGBColor(0x33, 0x33, 0x33),
        align=PP_ALIGN.CENTER)

    # ── 2. Agenda slide ───────────────────────────────────────────────────────
    if len(slides) > 1:
        ag = prs.slides.add_slide(blank)
        _bg(ag, WH)
        _header_strip(ag, "Agenda", Y, BK, logo)
        y = Inches(1.6)
        for idx, s in enumerate(slides, 1):
            # alternating row shading
            if idx % 2 == 0:
                _rect(ag, Inches(0.45), y - Inches(0.06),
                      Inches(12.44), Inches(0.5), RGBColor(0xF9, 0xF9, 0xF9))
            # yellow number badge
            _rect(ag, Inches(0.55), y + Inches(0.04), Inches(0.38), Inches(0.38), Y)
            _tb(ag, Inches(0.55), y + Inches(0.04), Inches(0.38), Inches(0.38),
                str(idx), Pt(12), BK, bold=True, align=PP_ALIGN.CENTER)
            _tb(ag, Inches(1.1), y, Inches(10.5), Inches(0.44),
                s.get("title", ""), Pt(14), DK)
            y += Inches(0.54)
        _footer(ag, Y, BK, GR, logo)

    # ── 3+. Content slides ────────────────────────────────────────────────────
    for slide_info in slides:
        s_title     = slide_info.get("title", "")
        bullets     = slide_info.get("bullets", [])
        is_takeaway = "takeaway" in s_title.lower() or "key" in s_title.lower()

        if is_takeaway:
            s = prs.slides.add_slide(blank)
            _bg(s, Y)
            _rect(s, 0, 0, Inches(0.22), Inches(7.5), BK)
            _tb(s, Inches(0.55), Inches(0.3), Inches(11), Inches(0.9),
                s_title, Pt(28), BK, bold=True)
            _rect(s, Inches(0.55), Inches(1.28), Inches(5), Inches(0.05), BK)
            y = Inches(1.5)
            for bullet in bullets[:6]:
                _rect(s, Inches(0.55), y + Inches(0.12), Inches(0.05), Inches(0.3), BK)
                _tb(s, Inches(0.78), y, Inches(11.5), Inches(0.62), bullet, Pt(16), BK)
                y += Inches(0.72)
            _rect(s, 0, Inches(6.9), Inches(13.33), Inches(0.6), RGBColor(0x11, 0x11, 0x11))
            _tb(s, Inches(0.55), Inches(6.93), Inches(12), Inches(0.45),
                "MTN Nigeria  ·  Confidential", Pt(9), GR)
            if logo:
                try:
                    s.shapes.add_picture(logo, Inches(11.8), Inches(6.85), Inches(1.25), Inches(0.6))
                except Exception:
                    pass
        else:
            s = prs.slides.add_slide(blank)
            _bg(s, WH)
            _header_strip(s, s_title, Y, BK, logo)
            y = Inches(1.6)
            for i, bullet in enumerate(bullets[:6]):
                # alternating subtle row tint
                if i % 2 == 0:
                    _rect(s, Inches(0.45), y - Inches(0.06),
                          Inches(12.44), Inches(0.66), RGBColor(0xFC, 0xFC, 0xFC))
                _rect(s, Inches(0.5), y + Inches(0.14), Inches(0.07), Inches(0.3), Y)
                _tb(s, Inches(0.75), y, Inches(11.5), Inches(0.64), bullet, Pt(15), DK)
                y += Inches(0.72)
            _footer(s, Y, BK, GR, logo)

    # ── Network heatmap slide ─────────────────────────────────────────────────
    if map_data:
        try:
            ms = prs.slides.add_slide(blank)
            _bg(ms, WH)
            _header_strip(ms, "Network Status — Nigeria Heatmap", Y, BK, logo)
            img_bytes = _generate_map_image(map_data)
            img_buf   = io.BytesIO(img_bytes)
            ms.shapes.add_picture(img_buf, Inches(0.35), Inches(1.38), Inches(12.63), Inches(5.45))
            _footer(ms, Y, BK, GR, logo)
        except Exception as _e:
            logger.warning(f"Map slide skipped: {_e}")

    # ── Fraud Intelligence slide ──────────────────────────────────────────────
    if fraud_data:
        try:
            RED = RGBColor(0xEF, 0x44, 0x44)
            AMB = RGBColor(0xF5, 0x9E, 0x0B)
            GRN = RGBColor(0x22, 0xC5, 0x5E)
            fs = prs.slides.add_slide(blank)
            _bg(fs, BK)
            _header_strip(fs, "Fraud Intelligence Summary", Y, BK, logo)

            hi  = fraud_data.get("high_risk", 0)
            med = fraud_data.get("medium_risk", 0)
            exp = fraud_data.get("total_exposure_ngn", 0)

            # Summary bar
            _tb(fs, Inches(0.5), Inches(1.38), Inches(12), Inches(0.5),
                f"{hi} HIGH RISK  ·  {med} MEDIUM  ·  Exposure: ₦{exp:,.0f}",
                Pt(14), Y, bold=True)

            _RISK_COLOR = {"HIGH": RED, "MEDIUM": AMB, "LOW": GRN}
            y = Inches(2.05)
            for s in (fraud_data.get("signals") or [])[:5]:
                risk_label = s.get("risk", "MEDIUM")
                clr = _RISK_COLOR.get(risk_label, AMB)
                _rect(fs, Inches(0.5), y + Inches(0.1), Inches(0.07), Inches(0.38), clr)
                _tb(fs, Inches(0.72), y, Inches(9.0), Inches(0.5),
                    f"[{risk_label}] {s.get('id', '')} — {s.get('title', '')}", Pt(11), WH, bold=True)
                detail = s.get("detail", "")[:100] + ("…" if len(s.get("detail", "")) > 100 else "")
                _tb(fs, Inches(0.72), y + Inches(0.46), Inches(11.5), Inches(0.36),
                    detail, Pt(9), GR)
                if s.get("amount_ngn"):
                    _tb(fs, Inches(9.9), y, Inches(3.0), Inches(0.5),
                        f"₦{s['amount_ngn']:,.0f}", Pt(12), Y, bold=True)
                y += Inches(1.0)

            _footer(fs, Y, BK, GR, logo)
        except Exception as _e:
            logger.warning(f"Fraud slide skipped: {_e}")

    # ── Sources slide ─────────────────────────────────────────────────────────
    if citations:
        cs = prs.slides.add_slide(blank)
        _bg(cs, LG)
        _header_strip(cs, "Sources & References", Y, BK, logo)
        y = Inches(1.6)
        for i, c in enumerate(citations[:9], 1):
            if isinstance(c, str):
                doc_title = c
            else:
                doc_title = (
                    c.get("document_title") or c.get("source")
                    or c.get("document_id", "Unknown")
                )
            _rect(cs, Inches(0.5), y + Inches(0.12), Inches(0.06), Inches(0.28), Y)
            _tb(cs, Inches(0.75), y, Inches(11.5), Inches(0.52),
                f"[{i}]  {doc_title}", Pt(13), DK)
            y += Inches(0.58)
        _footer(cs, Y, BK, GR, logo)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Slide helpers ─────────────────────────────────────────────────────────────

def _header_strip(slide, slide_title: str, Y, BK, logo=None):
    from pptx.util import Inches, Pt
    _rect(slide, 0, 0, Inches(13.33), Inches(1.3), Y)
    _rect(slide, 0, Inches(1.3), Inches(13.33), Inches(0.05), BK)
    _tb(slide, Inches(0.45), Inches(0.22), Inches(10.5), Inches(0.88),
        slide_title, Pt(23), BK, bold=True)
    if logo:
        try:
            slide.shapes.add_picture(logo, Inches(11.5), Inches(0.1), Inches(1.6), Inches(1.05))
        except Exception:
            pass


def _footer(slide, Y, BK, GR, logo=None):
    from pptx.util import Inches, Pt
    _rect(slide, 0, Inches(6.9), Inches(13.33), Inches(0.6), BK)
    _rect(slide, 0, Inches(6.9), Inches(0.22), Inches(0.6), Y)
    _tb(slide, Inches(0.45), Inches(6.93), Inches(9), Inches(0.45),
        "MTN Nigeria  ·  Enterprise Intelligence  ·  Confidential", Pt(8), GR)
    if logo:
        try:
            slide.shapes.add_picture(logo, Inches(11.8), Inches(6.85), Inches(1.3), Inches(0.6))
        except Exception:
            pass


def _bg(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _tb(slide, left, top, width, height, text: str, size,
        color, bold: bool = False, align=None, font_name: str = FONT_BODY):
    """Add a textbox with Calibri font. Guards against empty-text IndexError."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p     = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    if p.runs:
        run               = p.runs[0]
        run.font.size     = size
        run.font.bold     = bold
        run.font.name     = font_name
        run.font.color.rgb = color
    return txBox


# ── LLM slide structuring ─────────────────────────────────────────────────────

async def _structure_slides(title: str, answer: str) -> dict:
    try:
        from agents.kernel import llm_complete
        prompt = f"""Convert this intelligence report into a slide deck outline.
Return ONLY valid JSON — no markdown fences, no commentary.

Report title: "{title}"

Report content:
{answer[:3500]}

Rules:
- 4 to 6 content slides
- Each slide: concise punchy title (≤7 words) + 3–5 bullet points (≤18 words each)
- Preserve ALL key numbers, percentages, and named entities exactly as written
- Final slide MUST be titled "Key Takeaways" with 3–4 action-oriented bullets

JSON schema (return nothing else):
{{"slides": [{{"title": "...", "bullets": ["...", "..."]}}]}}"""

        raw   = await llm_complete(prompt, max_tokens=1200, temperature=0.2)
        clean = raw.strip().replace("```json", "").replace("```", "").strip()
        return json.loads(clean)
    except Exception as e:
        logger.warning(f"Slide structuring LLM call failed, using fallback: {e}")
        return _fallback_slides(title, answer)


def _fallback_slides(title: str, answer: str) -> dict:
    slides, current = [], None
    for line in answer.split("\n"):
        line = line.strip()
        is_heading = (
            line.startswith("## ") or line.startswith("### ")
            or (line.startswith("**") and line.endswith("**") and len(line) < 80)
        )
        if is_heading:
            if current and current["bullets"]:
                slides.append(current)
            heading = line.lstrip("#").strip().strip("*")
            current = {"title": heading[:60], "bullets": []}
        elif (line.startswith("- ") or line.startswith("* ")) and current is not None:
            bullet = re.sub(r"\*\*(.+?)\*\*", r"\1", line[2:])[:130]
            current["bullets"].append(bullet)
        elif line and not line.startswith("|") and current is not None and len(current["bullets"]) < 5:
            clean = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
            if len(clean) > 20:
                current["bullets"].append(clean[:130])
    if current and current["bullets"]:
        slides.append(current)
    return {"slides": slides[:6] or [{"title": title, "bullets": [answer[:200]]}]}
